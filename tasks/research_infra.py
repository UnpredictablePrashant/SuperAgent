import base64
import csv
import importlib.util
import json
import logging
import math
import mimetypes
import os
import re
import shutil
import subprocess
import tempfile
import threading
import time
import zipfile
from concurrent.futures import ThreadPoolExecutor, as_completed
from collections import deque
from pathlib import Path
from urllib.parse import urlencode, urljoin, urlparse
from urllib.request import Request, urlopen
import xml.etree.ElementTree as ET

from kendr.llm_router import (
    PROVIDER_CUSTOM,
    PROVIDER_GLM,
    PROVIDER_MINIMAX,
    PROVIDER_OLLAMA,
    PROVIDER_OPENAI,
    PROVIDER_OPENROUTER,
    PROVIDER_QWEN,
    PROVIDER_XAI,
    get_api_key,
    get_base_url,
    get_model_capabilities,
)
from tasks.utils import llm, model_selection_for_agent, normalize_llm_text

logger = logging.getLogger(__name__)

_DDGS_IMPORT_ERROR = ""
try:
    from ddgs import DDGS as _DDGS_CLIENT_FACTORY
except Exception as exc:
    _DDGS_IMPORT_ERROR = str(exc)
    try:
        from duckduckgo_search import DDGS as _DDGS_CLIENT_FACTORY
        _DDGS_IMPORT_ERROR = ""
    except Exception as compat_exc:
        _DDGS_CLIENT_FACTORY = None
        _DDGS_IMPORT_ERROR = str(compat_exc)

try:
    from ddgs.exceptions import RatelimitException as _DDGS_RATELIMIT_EXCEPTION
except Exception:
    try:
        from duckduckgo_search.exceptions import RatelimitException as _DDGS_RATELIMIT_EXCEPTION
    except Exception:
        class _DDGS_RATELIMIT_EXCEPTION(Exception):
            pass


DEFAULT_EMBEDDING_MODEL = os.getenv("OPENAI_EMBEDDING_MODEL", "text-embedding-3-small")
DEFAULT_VISION_MODEL = os.getenv(
    "OPENAI_VISION_MODEL",
    os.getenv("OPENAI_MODEL_GENERAL", os.getenv("OPENAI_MODEL", "gpt-4o-mini")),
)
DEFAULT_QDRANT_URL = os.getenv("QDRANT_URL", "http://localhost:6333")
DEFAULT_QDRANT_COLLECTION = os.getenv("QDRANT_COLLECTION", "research_memory")
SERP_API_URL = "https://serpapi.com/search.json"
OPENALEX_API_URL = "https://api.openalex.org/works"
ARXIV_API_URL = "https://export.arxiv.org/api/query"
REDDIT_BASE_URL = "https://www.reddit.com"
DUCKDUCKGO_INSTANT_ANSWER_URL = "https://api.duckduckgo.com/"
DUCKDUCKGO_HTML_SEARCH_URL = "https://html.duckduckgo.com/html/"
BROWSER_USE_SEARCH_ENGINES = {
    "google": "https://www.google.com/search",
    "duckduckgo": "https://duckduckgo.com/html/",
}
DDGS_MAX_RESULTS_PER_CALL = max(1, min(int(os.getenv("KENDR_DDGS_MAX_RESULTS_PER_CALL", "10") or 10), 10))
DDGS_MAX_CALLS_PER_REQUEST = max(1, min(int(os.getenv("KENDR_DDGS_MAX_CALLS_PER_REQUEST", "6") or 6), 10))
DDGS_RATE_LIMIT_RETRIES = max(1, min(int(os.getenv("KENDR_DDGS_RATE_LIMIT_RETRIES", "3") or 3), 6))
DDGS_RATE_LIMIT_BACKOFF = max(1.0, float(os.getenv("KENDR_DDGS_RATE_LIMIT_BACKOFF", "3.0") or 3.0))
DDGS_INTER_CALL_DELAY_SECONDS = max(0.0, float(os.getenv("KENDR_DDGS_INTER_CALL_DELAY_SECONDS", "1.0") or 1.0))
DDGS_LARGE_BATCH_EXTRA_DELAY_SECONDS = max(0.0, float(os.getenv("KENDR_DDGS_LARGE_BATCH_EXTRA_DELAY_SECONDS", "2.0") or 2.0))
DDGS_RECENT_MARKERS = {
    "today", "latest", "recent", "current", "new", "news", "breaking", "this week", "this month",
}
DDGS_ACADEMIC_MARKERS = {
    "academic", "paper", "papers", "study", "studies", "journal", "peer-reviewed", "scientific",
    "scholarly", "literature", "systematic review", "meta-analysis", "pubmed", "arxiv",
}
DDGS_PATENT_MARKERS = {
    "patent", "patents", "prior art", "prior-art", "inventor", "assignee", "claims", "uspto",
    "espacenet", "wipo",
}
SEARCH_BACKEND_ALIAS_MAP = {
    "": "auto",
    "auto": "auto",
    "duckduckgo": "duckduckgo",
    "ddgs": "duckduckgo",
    "duckduckgo_ddgs": "duckduckgo",
    "serpapi": "serpapi",
    "browser_use": "browser_use_mcp",
    "browser-use": "browser_use_mcp",
    "browser_use_mcp": "browser_use_mcp",
    "playwright": "playwright_browser",
    "playwright_browser": "playwright_browser",
}
IMAGE_FILE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tif", ".tiff"}
TEXT_LIKE_SOURCE_EXTENSIONS = {
    ".txt",
    ".md",
    ".json",
    ".html",
    ".htm",
    ".xml",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".cfg",
    ".conf",
    ".env",
    ".properties",
    ".sql",
    ".py",
    ".js",
    ".jsx",
    ".ts",
    ".tsx",
    ".java",
    ".go",
    ".rs",
    ".c",
    ".cc",
    ".cpp",
    ".h",
    ".hpp",
    ".cs",
    ".rb",
    ".php",
    ".swift",
    ".kt",
    ".kts",
    ".scala",
    ".sh",
    ".bash",
    ".zsh",
    ".ps1",
    ".bat",
    ".tf",
    ".tfvars",
    ".gradle",
}
LOCAL_DRIVE_SUPPORTED_EXTENSIONS = {
    *TEXT_LIKE_SOURCE_EXTENSIONS,
    ".csv",
    ".pdf",
    ".doc",
    ".docx",
    ".xls",
    ".xlsx",
    ".xlsm",
    ".ppt",
    ".pptx",
    ".pptm",
    *IMAGE_FILE_EXTENSIONS,
}

OPENAI_COMPAT_VISION_PROVIDERS = {
    PROVIDER_OPENAI,
    PROVIDER_XAI,
    PROVIDER_MINIMAX,
    PROVIDER_QWEN,
    PROVIDER_GLM,
    PROVIDER_OLLAMA,
    PROVIDER_OPENROUTER,
    PROVIDER_CUSTOM,
}


def strip_code_fences(text: str) -> str:
    stripped = (text or "").strip()
    if stripped.startswith("```") and stripped.endswith("```"):
        lines = stripped.splitlines()
        if len(lines) >= 2:
            return "\n".join(lines[1:-1]).strip()
    return stripped


_LLM_MAX_RETRIES = int(os.getenv("KENDR_LLM_MAX_RETRIES", "3"))
_LLM_BASE_DELAY = float(os.getenv("KENDR_LLM_BASE_DELAY", "2.0"))
_LLM_MAX_DELAY = float(os.getenv("KENDR_LLM_MAX_DELAY", "30.0"))

_TRANSIENT_ERROR_MARKERS = (
    "timeout", "timed out", "connection", "connect", "reset by peer",
    "temporarily unavailable", "rate limit", "429", "502", "503", "504",
    "eof", "broken pipe", "network", "ssl", "handshake",
)


def _is_transient_error(exc: Exception) -> bool:
    msg = str(exc).lower()
    return any(marker in msg for marker in _TRANSIENT_ERROR_MARKERS)


def _sanitize_unpaired_surrogates(text: str) -> str:
    if not text:
        return ""
    return re.sub(r"[\ud800-\udfff]", "\uFFFD", text)


def llm_text(prompt: str, *, max_retries: int = _LLM_MAX_RETRIES) -> str:
    safe_prompt = _sanitize_unpaired_surrogates(prompt)
    if safe_prompt != prompt:
        logger.warning("llm_text sanitized unpaired surrogate characters from prompt before LLM call.")
    last_exc: Exception | None = None
    for attempt in range(1, max_retries + 1):
        try:
            response = llm.invoke(safe_prompt)
            return normalize_llm_text(response.content if hasattr(response, "content") else response).strip()
        except Exception as exc:
            last_exc = exc
            if attempt >= max_retries or not _is_transient_error(exc):
                raise
            delay = min(_LLM_BASE_DELAY * (2 ** (attempt - 1)), _LLM_MAX_DELAY)
            logger.warning(
                "llm_text transient error (attempt %d/%d): %s — retrying in %.1fs",
                attempt, max_retries, exc, delay,
            )
            time.sleep(delay)
    raise last_exc  # unreachable but satisfies type checkers


def llm_json(prompt: str, fallback):
    try:
        return json.loads(strip_code_fences(llm_text(prompt)))
    except Exception:
        return fallback


def get_openai_client():
    from openai import OpenAI

    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        raise ValueError("OPENAI_API_KEY is required for this agent.")
    return OpenAI(api_key=api_key)


def _build_openai_compatible_client(provider: str):
    from openai import OpenAI

    normalized_provider = str(provider or "").strip().lower() or PROVIDER_OPENAI
    api_key = get_api_key(normalized_provider)
    if not api_key and normalized_provider in {PROVIDER_OLLAMA, PROVIDER_CUSTOM}:
        api_key = "ollama"
    if not api_key:
        raise ValueError(f"{normalized_provider.upper()} API key is required for OCR/image analysis.")

    init = {"api_key": api_key}
    base_url = str(get_base_url(normalized_provider) or "").strip()
    if base_url:
        init["base_url"] = base_url
    return OpenAI(**init)


def _resolve_vision_backend(
    *,
    provider: str = "",
    model: str = "",
    agent_name: str = "ocr_agent",
) -> dict[str, object]:
    explicit_provider = str(provider or "").strip().lower()
    explicit_model = str(model or "").strip()
    if explicit_provider and explicit_model:
        selection = {"provider": explicit_provider, "model": explicit_model, "source": "explicit"}
    else:
        selection = model_selection_for_agent(agent_name)

    selected_provider = str(selection.get("provider") or "").strip().lower() or PROVIDER_OPENAI
    selected_model = str(selection.get("model") or "").strip()
    selected_source = str(selection.get("source") or "").strip() or "agent_selection"
    trusted_selection = bool(explicit_provider and explicit_model) or selected_source == "runtime_override"
    capabilities = get_model_capabilities(selected_model, selected_provider)

    if (
        selected_provider in OPENAI_COMPAT_VISION_PROVIDERS
        and selected_model
        and (trusted_selection or bool(capabilities.get("vision")))
    ):
        return {
            "provider": selected_provider,
            "model": selected_model,
            "source": selected_source,
            "fallback": False,
        }

    fallback_model = str(DEFAULT_VISION_MODEL or "").strip()
    if fallback_model and str(get_api_key(PROVIDER_OPENAI) or "").strip():
        return {
            "provider": PROVIDER_OPENAI,
            "model": fallback_model,
            "source": "openai_fallback",
            "fallback": True,
        }

    if selected_provider in OPENAI_COMPAT_VISION_PROVIDERS and selected_model:
        return {
            "provider": selected_provider,
            "model": selected_model,
            "source": selected_source,
            "fallback": True,
        }

    raise ValueError(
        "No OCR-capable vision backend is available. Configure OPENAI_API_KEY or enable a supported vision model."
    )


def _image_data_uri(path: Path) -> str:
    mime_type = mimetypes.guess_type(path.name)[0] or "image/png"
    image_b64 = base64.b64encode(path.read_bytes()).decode("ascii")
    return f"data:{mime_type};base64,{image_b64}"


def _extract_chat_completion_text(response) -> str:
    output_text = getattr(response, "output_text", "") or ""
    if output_text:
        return str(output_text).strip()

    choices = getattr(response, "choices", None) or []
    for choice in choices:
        message = getattr(choice, "message", None)
        content = getattr(message, "content", "")
        if isinstance(content, str):
            if content.strip():
                return content.strip()
            continue
        if isinstance(content, list):
            parts: list[str] = []
            for item in content:
                if isinstance(item, dict):
                    text = item.get("text")
                    if isinstance(text, str) and text.strip():
                        parts.append(text.strip())
                else:
                    text = getattr(item, "text", "")
                    if isinstance(text, str) and text.strip():
                        parts.append(text.strip())
            if parts:
                return "\n".join(parts).strip()
    return ""


def _run_openai_compat_image_prompt(
    path_str: str,
    instruction: str,
    *,
    agent_name: str,
    provider: str = "",
    model: str = "",
) -> dict:
    backend = _resolve_vision_backend(provider=provider, model=model, agent_name=agent_name)
    path = Path(path_str)
    client = _build_openai_compatible_client(str(backend.get("provider") or ""))
    image_url = _image_data_uri(path)
    response = None
    try:
        response = client.chat.completions.create(
            model=str(backend.get("model") or ""),
            temperature=0,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": instruction},
                        {"type": "image_url", "image_url": {"url": image_url}},
                    ],
                }
            ],
        )
        output_text = _extract_chat_completion_text(response)
    except Exception:
        if str(backend.get("provider") or "").strip().lower() != PROVIDER_OPENAI:
            raise
        response = client.responses.create(
            model=str(backend.get("model") or ""),
            input=[
                {
                    "role": "user",
                    "content": [
                        {"type": "input_text", "text": instruction},
                        {"type": "input_image", "image_url": image_url},
                    ],
                }
            ],
        )
        output_text = getattr(response, "output_text", "") or ""
    payload = response.model_dump() if hasattr(response, "model_dump") else {"response": str(response)}
    return {
        "path": str(path),
        "raw": payload,
        "provider": str(backend.get("provider") or ""),
        "model": str(backend.get("model") or ""),
        "selection_source": str(backend.get("source") or ""),
        "text": str(output_text or "").strip(),
    }


def chunk_text(text: str, chunk_size: int = 1200, overlap: int = 150) -> list[str]:
    cleaned = " ".join((text or "").split())
    if not cleaned:
        return []
    chunks = []
    start = 0
    while start < len(cleaned):
        end = min(len(cleaned), start + chunk_size)
        chunks.append(cleaned[start:end])
        if end >= len(cleaned):
            break
        start = max(0, end - overlap)
    return chunks


def _regex_html_to_text(html: str) -> str:
    text = re.sub(r"(?is)<(script|style).*?>.*?</\\1>", " ", html)
    text = re.sub(r"(?s)<[^>]+>", " ", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def html_to_text(html: str) -> str:
    try:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        return re.sub(r"\s+", " ", soup.get_text(" ")).strip()
    except Exception:
        return _regex_html_to_text(html)


def fetch_url_content(url: str, timeout: int = 20) -> dict:
    if _browser_use_view_enabled():
        try:
            browser_payload = browser_use_fetch_url(url, timeout=timeout)
            if str(browser_payload.get("text", "") or "").strip():
                return browser_payload
        except Exception:
            pass

    request = Request(
        url,
        headers={
            "User-Agent": os.getenv(
                "RESEARCH_USER_AGENT",
                "multi-agent-research-bot/1.0 (+https://localhost)",
            )
        },
    )
    with urlopen(request, timeout=timeout) as response:
        body = response.read()
        content_type = response.headers.get("Content-Type", "")
    text = body.decode("utf-8", errors="ignore")
    return {
        "url": url,
        "content_type": content_type,
        "text": html_to_text(text) if "html" in content_type else text,
        "raw_text": text,
    }


def fetch_urls_content(
    urls: list[str],
    *,
    timeout: int = 20,
    max_workers: int | None = None,
    per_domain_limit: int | None = None,
    progress_callback=None,
) -> list[dict]:
    ordered_urls = [str(url or "").strip() for url in urls if str(url or "").strip()]
    if not ordered_urls:
        return []

    if len(ordered_urls) == 1:
        url = ordered_urls[0]
        if progress_callback:
            try:
                progress_callback(url, "started", None, 1, 1)
            except Exception:
                pass
        try:
            payload = dict(fetch_url_content(url, timeout=timeout) or {})
            payload["url"] = str(payload.get("url", "") or url).strip()
            payload["content_type"] = str(payload.get("content_type", "") or "").strip()
            payload["text"] = str(payload.get("text", "") or "")
            payload["raw_text"] = str(payload.get("raw_text", payload["text"]) or "")
            if progress_callback:
                try:
                    progress_callback(payload["url"], "completed", payload, 1, 1)
                except Exception:
                    pass
            return [payload]
        except Exception as exc:
            payload = {"url": url, "content_type": "", "text": "", "raw_text": "", "error": str(exc)}
            if progress_callback:
                try:
                    progress_callback(url, "failed", payload, 1, 1)
                except Exception:
                    pass
            return [payload]

    worker_count = max_workers if max_workers is not None else _link_worker_count(len(ordered_urls))
    worker_count = max(1, min(len(ordered_urls), int(worker_count)))
    domain_limit = max(1, min(worker_count, per_domain_limit if per_domain_limit is not None else _per_domain_link_limit()))
    results: list[dict | None] = [None] * len(ordered_urls)
    domain_semaphores: dict[str, threading.BoundedSemaphore] = {}
    domain_lock = threading.Lock()

    def _domain_semaphore(url: str) -> threading.BoundedSemaphore:
        domain = urlparse(url).netloc.lower() or "__default__"
        with domain_lock:
            semaphore = domain_semaphores.get(domain)
            if semaphore is None:
                semaphore = threading.BoundedSemaphore(domain_limit)
                domain_semaphores[domain] = semaphore
            return semaphore

    def _fetch_one(index: int, url: str) -> tuple[int, dict]:
        semaphore = _domain_semaphore(url)
        with semaphore:
            if progress_callback:
                try:
                    progress_callback(url, "started", None, index + 1, len(ordered_urls))
                except Exception:
                    pass
            try:
                payload = fetch_url_content(url, timeout=timeout)
                normalized = dict(payload or {})
                normalized["url"] = str(normalized.get("url", "") or url).strip()
                normalized["content_type"] = str(normalized.get("content_type", "") or "").strip()
                normalized["text"] = str(normalized.get("text", "") or "")
                normalized["raw_text"] = str(normalized.get("raw_text", normalized["text"]) or "")
                if progress_callback:
                    try:
                        progress_callback(normalized["url"], "completed", normalized, index + 1, len(ordered_urls))
                    except Exception:
                        pass
                return index, normalized
            except Exception as exc:
                payload = {"url": url, "content_type": "", "text": "", "raw_text": "", "error": str(exc)}
                if progress_callback:
                    try:
                        progress_callback(url, "failed", payload, index + 1, len(ordered_urls))
                    except Exception:
                        pass
                return index, payload

    with ThreadPoolExecutor(max_workers=worker_count) as pool:
        futures = [pool.submit(_fetch_one, index, url) for index, url in enumerate(ordered_urls)]
        for future in as_completed(futures):
            index, payload = future.result()
            results[index] = payload

    return [item for item in results if isinstance(item, dict)]


def _env_flag(name: str, default: bool = True) -> bool:
    raw = str(os.getenv(name, "1" if default else "0") or "").strip().lower()
    if raw in {"", "1", "true", "yes", "on"}:
        return True
    if raw in {"0", "false", "no", "off"}:
        return False
    return default


def _env_int(name: str, default: int, *, minimum: int = 1, maximum: int | None = None) -> int:
    raw = str(os.getenv(name, str(default)) or "").strip()
    try:
        value = int(raw)
    except Exception:
        value = default
    value = max(minimum, value)
    if maximum is not None:
        value = min(maximum, value)
    return value


def _bounded_worker_count(total_items: int, default: int, *, env_name: str, maximum: int = 16) -> int:
    if total_items <= 0:
        return 1
    configured = _env_int(env_name, default, minimum=1, maximum=maximum)
    return max(1, min(total_items, configured))


def _link_worker_count(total_items: int) -> int:
    return _bounded_worker_count(total_items, 4, env_name="KENDR_RESEARCH_LINK_MAX_WORKERS", maximum=16)


def _per_domain_link_limit() -> int:
    return _env_int("KENDR_RESEARCH_PER_DOMAIN_MAX_WORKERS", 2, minimum=1, maximum=8)


def _document_worker_count(total_items: int) -> int:
    return _bounded_worker_count(total_items, 4, env_name="KENDR_RESEARCH_DOCUMENT_MAX_WORKERS", maximum=12)


def _heavy_document_worker_count() -> int:
    return _env_int("KENDR_RESEARCH_HEAVY_DOCUMENT_MAX_WORKERS", 2, minimum=1, maximum=6)


def _ocr_document_worker_count() -> int:
    return _env_int("KENDR_RESEARCH_OCR_MAX_WORKERS", 1, minimum=1, maximum=4)


def _browser_use_search_enabled() -> bool:
    return _env_flag("KENDR_BROWSER_USE_SEARCH_ENABLED", True)


def _browser_use_view_enabled() -> bool:
    return _env_flag("KENDR_BROWSER_USE_VIEW_ENABLED", True)


def _browser_use_search_engine() -> str:
    engine = str(os.getenv("KENDR_BROWSER_USE_SEARCH_ENGINE", "google") or "google").strip().lower()
    return engine if engine in BROWSER_USE_SEARCH_ENGINES else "google"


def _browser_use_server_enabled() -> bool:
    if not _browser_use_search_enabled() and not _browser_use_view_enabled():
        return False
    try:
        from kendr.mcp_manager import browser_use_server

        return isinstance(browser_use_server(enabled_only=True), dict)
    except Exception:
        return False


def _browser_use_result_text(payload: dict) -> str:
    text = str(payload.get("text", "") or "").strip()
    if text:
        return text
    content = payload.get("content")
    if isinstance(content, list):
        return "\n".join(str(item) for item in content).strip()
    return str(content or "").strip()


def _browser_use_search_url(query: str) -> str:
    engine = _browser_use_search_engine()
    base = BROWSER_USE_SEARCH_ENGINES.get(engine, BROWSER_USE_SEARCH_ENGINES["google"])
    params = {"q": query}
    if engine == "duckduckgo":
        params["kl"] = "us-en"
    return f"{base}?{urlencode(params)}"


def _extract_browser_use_results(text: str, *, num: int, page_url: str) -> list[dict]:
    results: list[dict] = []
    seen: set[str] = set()
    markdown_links = re.findall(r"\[([^\]]+)\]\((https?://[^)]+)\)", text or "")
    plain_links = re.findall(r"(https?://[^\s)>\]]+)", text or "")

    for title, url in markdown_links:
        clean_url = _clean_search_result_url(url)
        clean_title = re.sub(r"\s+", " ", str(title or "").strip())
        if not clean_url or clean_url in seen:
            continue
        seen.add(clean_url)
        results.append({
            "title": clean_title or clean_url,
            "url": clean_url,
            "snippet": "",
            "source": "browser-use MCP",
            "date": "",
        })
        if len(results) >= num:
            return results

    for url in plain_links:
        clean_url = _clean_search_result_url(url)
        if not clean_url or clean_url in seen or clean_url == page_url:
            continue
        seen.add(clean_url)
        results.append({
            "title": clean_url,
            "url": clean_url,
            "snippet": "",
            "source": "browser-use MCP",
            "date": "",
        })
        if len(results) >= num:
            break

    return results[:num]


def browser_use_fetch_url(url: str, *, timeout: int = 20) -> dict:
    from kendr.mcp_manager import browser_use_server, call_tool

    server = browser_use_server(enabled_only=True)
    if not isinstance(server, dict):
        raise RuntimeError("browser-use MCP server is not enabled")
    server_id = str(server.get("id") or server.get("server_id") or "").strip()
    navigate = call_tool(server_id, "browser_navigate", {"url": url}, timeout=max(20, timeout))
    extracted = call_tool(server_id, "browser_extract_content", {}, timeout=max(20, timeout))
    text = _browser_use_result_text(extracted)
    return {
        "url": url,
        "content_type": "text/markdown",
        "text": text,
        "raw_text": text,
        "provider": "browser_use_mcp",
        "navigate_text": _browser_use_result_text(navigate),
    }


def browser_use_search(query: str, *, num: int = 10) -> dict:
    search_url = _browser_use_search_url(query)
    payload = browser_use_fetch_url(search_url, timeout=30)
    text = str(payload.get("text", "") or "").strip()
    results = _extract_browser_use_results(text, num=num, page_url=search_url)
    return {
        "results": results[:num],
        "raw_text": text,
        "search_url": search_url,
        "engine": _browser_use_search_engine(),
    }


def extract_links(base_url: str, html: str, limit: int = 20, same_domain: bool = True) -> list[str]:
    matches = re.findall(r'href=["\\\'](.*?)["\\\']', html or "", flags=re.IGNORECASE)
    links = []
    base_domain = urlparse(base_url).netloc
    for href in matches:
        absolute = urljoin(base_url, href)
        parsed = urlparse(absolute)
        if parsed.scheme not in {"http", "https"}:
            continue
        if same_domain and parsed.netloc != base_domain:
            continue
        if absolute not in links:
            links.append(absolute)
        if len(links) >= limit:
            break
    return links


def crawl_urls(urls: list[str], max_pages: int = 5, same_domain: bool = True) -> list[dict]:
    pending = deque(str(url or "").strip() for url in urls if str(url or "").strip())
    enqueued = set(pending)
    visited: set[str] = set()
    pages: list[dict] = []
    worker_count = _link_worker_count(max_pages)

    while pending and len(pages) < max_pages:
        batch: list[str] = []
        remaining = max_pages - len(pages)
        while pending and len(batch) < min(worker_count, remaining):
            url = pending.popleft()
            enqueued.discard(url)
            if url in visited:
                continue
            visited.add(url)
            batch.append(url)
        if not batch:
            continue
        fetched_pages = fetch_urls_content(batch, timeout=20, max_workers=min(worker_count, len(batch)))
        for page in fetched_pages:
            page_url = str(page.get("url", "")).strip()
            content_type = str(page.get("content_type", "")).strip()
            pages.append(
                {
                    "url": page_url,
                    "content_type": content_type,
                    "text": str(page.get("text", "") or "")[:12000],
                    **({"error": str(page.get("error", "")).strip()} if page.get("error") else {}),
                }
            )
            if len(pages) >= max_pages:
                break
            raw_text = str(page.get("raw_text", "") or "")
            if "html" in content_type and raw_text:
                for link in extract_links(page_url, raw_text, same_domain=same_domain):
                    if link not in visited and link not in enqueued:
                        pending.append(link)
                        enqueued.add(link)
    return pages


def serp_search(query: str, *, engine: str = "google", num: int = 5, extra_params: dict | None = None) -> dict:
    api_key = os.getenv("SERP_API_KEY")
    if not api_key:
        raise ValueError("SERP_API_KEY is required for search-based agents.")
    params = {
        "engine": engine,
        "q": query,
        "api_key": api_key,
        "num": num,
        "hl": "en",
        "gl": "us",
    }
    if extra_params:
        params.update({k: v for k, v in extra_params.items() if v is not None})
    request_url = f"{SERP_API_URL}?{urlencode(params)}"
    with urlopen(request_url, timeout=30) as response:
        return json.loads(response.read().decode("utf-8"))


def serp_scholar_search(query: str, *, num: int = 10, start: int = 0, extra_params: dict | None = None) -> dict:
    params = {"start": start}
    if extra_params:
        params.update(extra_params)
    return serp_search(query, engine="google_scholar", num=num, extra_params=params)


def serp_patent_search(
    query: str,
    *,
    num: int = 10,
    page: int = 1,
    include_scholar: bool = False,
    extra_params: dict | None = None,
) -> dict:
    params = {
        "page": page,
        "num": num,
        "patents": "true",
        "scholar": "true" if include_scholar else "false",
    }
    if extra_params:
        params.update(extra_params)
    return serp_search(query, engine="google_patents", num=num, extra_params=params)


def _clean_search_result_url(url: str) -> str:
    value = str(url or "").strip()
    if not value:
        return ""
    parsed = urlparse(value)
    if parsed.scheme not in {"http", "https"}:
        return ""
    return value


def _ddgs_available() -> bool:
    return _DDGS_CLIENT_FACTORY is not None


def _playwright_search_available() -> bool:
    return importlib.util.find_spec("playwright") is not None


def normalize_research_search_backend(value: str = "") -> str:
    normalized = str(value or "").strip().lower()
    return SEARCH_BACKEND_ALIAS_MAP.get(normalized, "auto")


def research_search_backend_statuses() -> list[dict]:
    serpapi_ready = bool(os.getenv("SERP_API_KEY", "").strip())
    browser_use_ready = _browser_use_server_enabled()
    ddgs_ready = _ddgs_available()
    playwright_ready = _playwright_search_available()
    return [
        {
            "id": "auto",
            "label": "Auto",
            "enabled": True,
            "authenticated": False,
            "rate_limited": False,
            "description": "Prefer the strongest configured backend, then fall back automatically.",
            "note": "Uses SerpAPI first when configured; otherwise falls back to local/free search helpers.",
            "warning": (
                ""
                if serpapi_ready or not ddgs_ready
                else (
                    "Auto will likely fall back to unauthenticated DDGS search in this setup. "
                    "That path can rate-limit during heavier runs, even though Kendr caps calls and adds delays."
                )
            ),
        },
        {
            "id": "duckduckgo",
            "label": "DuckDuckGo (DDGS)",
            "enabled": ddgs_ready,
            "authenticated": False,
            "rate_limited": True,
            "description": "Unauthenticated DuckDuckGo search helper with bounded retries and delays.",
            "note": "No API key required.",
            "warning": (
                "Unauthenticated DDGS search can hit rate limits on heavier runs. "
                "Kendr keeps DDGS calls capped and delayed, but retries can still slow collection."
            ),
        },
        {
            "id": "serpapi",
            "label": "SerpAPI",
            "enabled": serpapi_ready,
            "authenticated": True,
            "rate_limited": False,
            "description": "Structured authenticated search results via SerpAPI.",
            "note": "Requires SERP_API_KEY.",
            "warning": "",
        },
        {
            "id": "browser_use_mcp",
            "label": "Browser-Use MCP",
            "enabled": browser_use_ready,
            "authenticated": False,
            "rate_limited": False,
            "description": "Browser-backed search extraction through the browser-use MCP server.",
            "note": "Requires a running browser-use MCP server.",
            "warning": "",
        },
        {
            "id": "playwright_browser",
            "label": "Playwright Browser",
            "enabled": playwright_ready,
            "authenticated": False,
            "rate_limited": False,
            "description": "Headless browser fallback that extracts results directly from the search page.",
            "note": "Requires Playwright plus an installed browser runtime.",
            "warning": "",
        },
    ]


def _is_ddgs_rate_limit_error(exc: Exception) -> bool:
    if isinstance(exc, _DDGS_RATELIMIT_EXCEPTION):
        return True
    lowered = str(exc or "").strip().lower()
    return "rate" in lowered and "limit" in lowered or "429" in lowered


def _ddgs_with_retry(fn, *, retries: int = DDGS_RATE_LIMIT_RETRIES, backoff: float = DDGS_RATE_LIMIT_BACKOFF):
    last_exc: Exception | None = None
    for attempt in range(1, retries + 1):
        try:
            return fn()
        except Exception as exc:
            last_exc = exc
            if attempt >= retries or not _is_ddgs_rate_limit_error(exc):
                raise
            wait_seconds = min(45.0, backoff ** attempt)
            logger.warning(
                "DDGS rate limit encountered (attempt %d/%d): %s — retrying in %.1fs",
                attempt,
                retries,
                exc,
                wait_seconds,
            )
            time.sleep(wait_seconds)
    if last_exc is not None:
        raise last_exc
    raise RuntimeError("DDGS retry failed without an explicit exception.")


def _ddgs_call_budget(*, num: int, max_search_calls: int = 0) -> int:
    requested = int(max_search_calls or 0)
    if requested <= 0:
        if num <= 3:
            requested = 2
        elif num <= 7:
            requested = 3
        else:
            requested = 4
    return max(1, min(requested, DDGS_MAX_CALLS_PER_REQUEST))


def _ddgs_inter_call_delay(call_index: int, total_calls: int) -> float:
    if call_index >= total_calls:
        return 0.0
    delay = DDGS_INTER_CALL_DELAY_SECONDS
    if total_calls >= 4 and call_index >= 3:
        delay += DDGS_LARGE_BATCH_EXTRA_DELAY_SECONDS
    return delay


def _search_focus_terms(text: str, *, limit: int = 6) -> str:
    stop_words = {
        "the", "and", "for", "with", "from", "that", "this", "into", "about", "after", "before",
        "using", "used", "should", "would", "could", "these", "those", "their", "there", "which",
        "what", "when", "where", "while", "than", "then", "into", "onto", "your", "our", "brief",
        "focused", "focus", "report", "research", "summary", "sources", "evidence",
    }
    terms: list[str] = []
    seen: set[str] = set()
    for token in re.findall(r"[a-z0-9][a-z0-9+./:-]{1,}", str(text or "").lower()):
        if token in stop_words or len(token) < 3:
            continue
        if token in seen:
            continue
        seen.add(token)
        terms.append(token)
        if len(terms) >= limit:
            break
    return " ".join(terms)


def _infer_duckduckgo_scope(query: str, *, focused_brief: str = "", research_scope: str = "") -> str:
    explicit = str(research_scope or "").strip().lower()
    if explicit in {"academic", "academics", "scholar", "scholarly", "papers"}:
        return "academic"
    if explicit in {"patent", "patents"}:
        return "patents"
    if explicit in {"news", "recent"}:
        return "news"
    haystack = f"{query} {focused_brief}".lower()
    if any(marker in haystack for marker in DDGS_PATENT_MARKERS):
        return "patents"
    if any(marker in haystack for marker in DDGS_ACADEMIC_MARKERS):
        return "academic"
    if any(marker in haystack for marker in DDGS_RECENT_MARKERS):
        return "news"
    return "web"


def _ddgs_timelimit_for_query(query: str, *, scope: str) -> str | None:
    lowered = str(query or "").lower()
    if scope == "news":
        if "today" in lowered or "breaking" in lowered:
            return "d"
        if "this week" in lowered or "weekly" in lowered:
            return "w"
        return "m"
    if "latest" in lowered or "recent" in lowered or "current" in lowered:
        return "m"
    return None


def _build_duckduckgo_query_plan(
    query: str,
    *,
    focused_brief: str = "",
    research_scope: str = "",
    max_calls: int,
) -> tuple[str, list[dict[str, str]]]:
    base_query = re.sub(r"\s+", " ", str(query or "").strip())
    scope = _infer_duckduckgo_scope(base_query, focused_brief=focused_brief, research_scope=research_scope)
    focus_terms = _search_focus_terms(focused_brief)
    candidate_queries: list[str] = [base_query]
    if focus_terms and focus_terms.lower() not in base_query.lower():
        candidate_queries.append(f"{base_query} {focus_terms}")

    if scope == "academic":
        candidate_queries.extend(
            [
                f"{base_query} peer reviewed study OR systematic review",
                f"site:pubmed.ncbi.nlm.nih.gov {base_query}",
                f"site:arxiv.org {base_query}",
                f"site:openalex.org OR site:nih.gov {base_query}",
            ]
        )
    elif scope == "patents":
        candidate_queries.extend(
            [
                f"{base_query} patent prior art",
                f"site:patents.google.com {base_query}",
                f"site:worldwide.espacenet.com {base_query}",
                f"site:uspto.gov {base_query}",
            ]
        )
    elif scope == "news":
        candidate_queries.extend(
            [
                f"{base_query} latest",
                f"{base_query} official announcement",
                f"{base_query} press release",
            ]
        )

    plan: list[dict[str, str]] = []
    seen: set[str] = set()
    for item in candidate_queries:
        normalized = re.sub(r"\s+", " ", str(item or "").strip())
        if not normalized:
            continue
        lowered = normalized.lower()
        if lowered in seen:
            continue
        seen.add(lowered)
        plan.append(
            {
                "query": normalized,
                "mode": "text",
                "timelimit": _ddgs_timelimit_for_query(normalized, scope=scope) or "",
                "scope": scope,
            }
        )
        if len(plan) >= max_calls:
            break
    return scope, plan


def _normalize_ddgs_result(item: dict, *, fallback_source: str = "DuckDuckGo DDGS") -> dict | None:
    if not isinstance(item, dict):
        return None
    url = _clean_search_result_url(item.get("href") or item.get("url") or item.get("link") or item.get("content") or "")
    title = re.sub(r"\s+", " ", str(item.get("title") or item.get("heading") or "")).strip()
    if not url or not title:
        return None
    snippet = re.sub(
        r"\s+",
        " ",
        str(item.get("body") or item.get("snippet") or item.get("excerpt") or item.get("description") or ""),
    ).strip()
    return {
        "title": title,
        "url": url,
        "snippet": snippet,
        "source": str(item.get("source") or fallback_source).strip() or fallback_source,
        "date": str(item.get("date") or item.get("published") or "").strip(),
    }


def ddgs_text_search(query: str, *, max_results: int = 10, region: str = "us-en", timelimit: str | None = None) -> list[dict]:
    if not _ddgs_available():
        raise RuntimeError(f"ddgs package is unavailable: {_DDGS_IMPORT_ERROR or 'not installed'}")
    bounded_results = max(1, min(int(max_results or DDGS_MAX_RESULTS_PER_CALL), DDGS_MAX_RESULTS_PER_CALL))

    def _run() -> list[dict]:
        client = _DDGS_CLIENT_FACTORY()
        payload = client.text(
            query=query,
            region=region,
            safesearch="moderate",
            timelimit=timelimit,
            max_results=bounded_results,
        )
        return list(payload or [])

    return _ddgs_with_retry(_run)


def duckduckgo_instant_answer(query: str) -> dict:
    params = {
        "q": str(query or "").strip(),
        "format": "json",
        "no_redirect": "1",
        "no_html": "1",
        "skip_disambig": "1",
    }
    request = Request(
        f"{DUCKDUCKGO_INSTANT_ANSWER_URL}?{urlencode(params)}",
        headers={
            "User-Agent": os.getenv(
                "RESEARCH_USER_AGENT",
                "multi-agent-research-bot/1.0 (+https://localhost)",
            ),
        },
        method="GET",
    )
    with urlopen(request, timeout=20) as response:
        raw = json.loads(response.read().decode("utf-8", errors="ignore"))
    return {
        "heading": str(raw.get("Heading", "") or "").strip(),
        "abstract": str(raw.get("AbstractText", "") or "").strip(),
        "abstract_url": str(raw.get("AbstractURL", "") or "").strip(),
        "answer": str(raw.get("Answer", "") or "").strip(),
        "answer_type": str(raw.get("AnswerType", "") or "").strip(),
        "definition": str(raw.get("Definition", "") or "").strip(),
        "definition_url": str(raw.get("DefinitionURL", "") or "").strip(),
        "entity": str(raw.get("Entity", "") or "").strip(),
        "related_topics_count": len(raw.get("RelatedTopics", []) if isinstance(raw.get("RelatedTopics", []), list) else []),
    }


def duckduckgo_ddgs_search(
    query: str,
    *,
    num: int = 10,
    focused_brief: str = "",
    research_scope: str = "",
    max_search_calls: int = 0,
    include_instant_answer: bool = True,
) -> dict:
    if not _ddgs_available():
        raise RuntimeError(f"ddgs package is unavailable: {_DDGS_IMPORT_ERROR or 'not installed'}")

    bounded_num = max(1, min(int(num or DDGS_MAX_RESULTS_PER_CALL), DDGS_MAX_RESULTS_PER_CALL))
    call_budget = _ddgs_call_budget(num=bounded_num, max_search_calls=max_search_calls)
    reserve_instant_answer = bool(include_instant_answer and call_budget > 1)
    query_call_budget = max(1, call_budget - (1 if reserve_instant_answer else 0))
    scope, query_plan = _build_duckduckgo_query_plan(
        query,
        focused_brief=focused_brief,
        research_scope=research_scope,
        max_calls=query_call_budget,
    )
    per_call_results = min(
        DDGS_MAX_RESULTS_PER_CALL,
        max(3, int(math.ceil(bounded_num / max(1, len(query_plan)))) + 1),
    )
    combined_results: list[dict] = []
    raw_batches: list[dict] = []
    seen_urls: set[str] = set()
    minimum_calls = 2 if scope in {"academic", "patents"} and len(query_plan) > 1 else 1

    for index, plan_item in enumerate(query_plan, start=1):
        raw_results = ddgs_text_search(
            plan_item["query"],
            max_results=per_call_results,
            timelimit=plan_item["timelimit"] or None,
        )
        raw_batches.append(
            {
                "query": plan_item["query"],
                "mode": plan_item["mode"],
                "timelimit": plan_item["timelimit"],
                "results_count": len(raw_results),
                "results": raw_results,
            }
        )
        for item in raw_results:
            normalized = _normalize_ddgs_result(item)
            if not normalized:
                continue
            url = normalized["url"]
            if url in seen_urls:
                continue
            seen_urls.add(url)
            combined_results.append(normalized)
            if len(combined_results) >= bounded_num and index >= minimum_calls:
                break
        if len(combined_results) >= bounded_num and index >= minimum_calls:
            break
        wait_seconds = _ddgs_inter_call_delay(index, len(query_plan))
        if wait_seconds > 0:
            time.sleep(wait_seconds)

    instant_answer: dict = {}
    if reserve_instant_answer:
        try:
            instant_answer = duckduckgo_instant_answer(query)
        except Exception as exc:
            instant_answer = {"error": str(exc)}

    return {
        "results": combined_results[:bounded_num],
        "query_plan": query_plan,
        "raw_batches": raw_batches,
        "instant_answer": instant_answer,
        "scope": scope,
        "call_budget": call_budget,
    }


def duckduckgo_html_search(query: str, *, num: int = 10) -> dict:
    payload = urlencode({"q": query, "kl": "us-en"})
    request = Request(
        DUCKDUCKGO_HTML_SEARCH_URL,
        data=payload.encode("utf-8"),
        headers={
            "User-Agent": os.getenv(
                "RESEARCH_USER_AGENT",
                "multi-agent-research-bot/1.0 (+https://localhost)",
            ),
            "Content-Type": "application/x-www-form-urlencoded",
        },
        method="POST",
    )
    with urlopen(request, timeout=30) as response:
        html = response.read().decode("utf-8", errors="ignore")

    results: list[dict] = []
    try:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(html, "html.parser")
        anchors = soup.select("a.result__a, a.result-link")
        for anchor in anchors:
            title = re.sub(r"\s+", " ", anchor.get_text(" ", strip=True)).strip()
            href = _clean_search_result_url(anchor.get("href", ""))
            if not title or not href:
                continue
            container = anchor.find_parent(class_="result") or anchor.parent
            snippet_el = None
            if container is not None:
                snippet_el = container.select_one(".result__snippet, .result-snippet")
            snippet = ""
            if snippet_el is not None:
                snippet = re.sub(r"\s+", " ", snippet_el.get_text(" ", strip=True)).strip()
            if href not in {item.get("url", "") for item in results}:
                results.append(
                    {
                        "title": title,
                        "url": href,
                        "snippet": snippet,
                        "source": "DuckDuckGo",
                        "date": "",
                    }
                )
            if len(results) >= num:
                break
    except Exception:
        pass

    if not results:
        for match in re.finditer(
            r'<a[^>]+class="[^"]*(?:result__a|result-link)[^"]*"[^>]+href="([^"]+)"[^>]*>(.*?)</a>',
            html,
            flags=re.IGNORECASE | re.DOTALL,
        ):
            href = _clean_search_result_url(match.group(1))
            title = re.sub(r"\s+", " ", re.sub(r"<[^>]+>", " ", match.group(2))).strip()
            if not href or not title:
                continue
            if href not in {item.get("url", "") for item in results}:
                results.append(
                    {
                        "title": title,
                        "url": href,
                        "snippet": "",
                        "source": "DuckDuckGo",
                        "date": "",
                    }
                )
            if len(results) >= num:
                break

    return {"results": results[:num], "raw_html": html}


def browser_search(query: str, *, num: int = 10) -> dict:
    try:
        from playwright.sync_api import sync_playwright
    except Exception as exc:
        raise RuntimeError(f"Playwright unavailable: {exc}") from exc

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page()
        try:
            page.goto(
                f"https://duckduckgo.com/html/?{urlencode({'q': query, 'kl': 'us-en'})}",
                wait_until="domcontentloaded",
                timeout=30000,
            )
            page.wait_for_timeout(1200)
            html = page.content()
        finally:
            browser.close()

    results: list[dict] = []
    try:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(html, "html.parser")
        anchors = soup.select("a.result__a, a.result-link")
        for anchor in anchors:
            title = re.sub(r"\s+", " ", anchor.get_text(" ", strip=True)).strip()
            href = _clean_search_result_url(anchor.get("href", ""))
            if not title or not href:
                continue
            container = anchor.find_parent(class_="result") or anchor.parent
            snippet_el = None
            if container is not None:
                snippet_el = container.select_one(".result__snippet, .result-snippet")
            snippet = ""
            if snippet_el is not None:
                snippet = re.sub(r"\s+", " ", snippet_el.get_text(" ", strip=True)).strip()
            if href not in {item.get("url", "") for item in results}:
                results.append(
                    {
                        "title": title,
                        "url": href,
                        "snippet": snippet,
                        "source": "DuckDuckGo browser",
                        "date": "",
                    }
                )
            if len(results) >= num:
                break
    except Exception as exc:
        raise RuntimeError(f"Browser search parsing failed: {exc}") from exc

    return {"results": results[:num], "raw_html": html}


def fetch_search_results(
    query: str,
    *,
    num: int = 10,
    fetch_pages: int = 3,
    progress_callback=None,
    provider_hint: str = "",
    focused_brief: str = "",
    research_scope: str = "",
    max_search_calls: int = 0,
    include_instant_answer: bool = True,
) -> dict:
    attempts: list[str] = []
    search_results: list[dict] = []
    provider = ""
    raw_payload: dict | None = None
    errors: list[str] = []
    bounded_num = max(1, min(int(num or DDGS_MAX_RESULTS_PER_CALL), DDGS_MAX_RESULTS_PER_CALL))

    serpapi_enabled = bool(os.getenv("SERP_API_KEY", "").strip())
    browser_use_enabled = _browser_use_server_enabled()
    normalized_backend = normalize_research_search_backend(provider_hint)

    provider_order = {
        "auto": ["serpapi", "browser_use_mcp", "ddgs", "duckduckgo_html", "playwright_browser"],
        "duckduckgo": ["ddgs", "duckduckgo_html", "browser_use_mcp", "playwright_browser", "serpapi"],
        "serpapi": ["serpapi", "ddgs", "duckduckgo_html", "browser_use_mcp", "playwright_browser"],
        "browser_use_mcp": ["browser_use_mcp", "ddgs", "duckduckgo_html", "playwright_browser", "serpapi"],
        "playwright_browser": ["playwright_browser", "ddgs", "duckduckgo_html", "browser_use_mcp", "serpapi"],
    }.get(normalized_backend, ["serpapi", "browser_use_mcp", "ddgs", "duckduckgo_html", "playwright_browser"])

    def _run_provider(provider_key: str) -> tuple[str, dict, list[dict]]:
        if provider_key == "serpapi":
            if not serpapi_enabled:
                raise RuntimeError("SERP_API_KEY not configured")
            serp_payload = serp_search(query, num=bounded_num)
            normalized_results: list[dict] = []
            for item in serp_payload.get("organic_results", [])[:bounded_num]:
                url = _clean_search_result_url(item.get("link", ""))
                if not url:
                    continue
                normalized_results.append(
                    {
                        "title": str(item.get("title", "")).strip(),
                        "url": url,
                        "snippet": str(item.get("snippet", "")).strip(),
                        "source": str(item.get("source", "")).strip(),
                        "date": str(item.get("date", "")).strip(),
                    }
                )
            return "serpapi", serp_payload, normalized_results
        if provider_key == "browser_use_mcp":
            if not browser_use_enabled:
                raise RuntimeError("browser-use MCP search server is not enabled")
            browser_use_payload = browser_use_search(query, num=bounded_num)
            return "browser_use_mcp", browser_use_payload, list(browser_use_payload.get("results", []) or [])
        if provider_key == "ddgs":
            ddgs_payload = duckduckgo_ddgs_search(
                query,
                num=bounded_num,
                focused_brief=focused_brief,
                research_scope=research_scope,
                max_search_calls=max_search_calls,
                include_instant_answer=include_instant_answer,
            )
            return "ddgs", ddgs_payload, list(ddgs_payload.get("results", []) or [])
        if provider_key == "duckduckgo_html":
            ddg_payload = duckduckgo_html_search(query, num=bounded_num)
            return "duckduckgo_html", ddg_payload, list(ddg_payload.get("results", []) or [])
        if provider_key == "playwright_browser":
            if not _playwright_search_available():
                raise RuntimeError("Playwright is not installed")
            browser_payload = browser_search(query, num=bounded_num)
            return "playwright_browser", browser_payload, list(browser_payload.get("results", []) or [])
        raise RuntimeError(f"Unknown search provider: {provider_key}")

    for provider_key in provider_order:
        if search_results:
            break
        if provider_key == "serpapi" and not serpapi_enabled:
            continue
        if provider_key == "browser_use_mcp" and not browser_use_enabled:
            continue
        if provider_key == "playwright_browser" and not _playwright_search_available():
            continue
        attempts.append(provider_key)
        try:
            provider, raw_payload, search_results = _run_provider(provider_key)
        except Exception as exc:
            errors.append(f"{provider_key}: {exc}")

    viewed_pages: list[dict] = []
    if search_results:
        view_urls = [
            _clean_search_result_url(item.get("url", ""))
            for item in search_results[: max(0, int(fetch_pages or 0))]
        ]
        view_urls = [url for url in view_urls if url]
        for page in fetch_urls_content(
            view_urls,
            timeout=20,
            max_workers=min(_link_worker_count(len(view_urls)), len(view_urls) or 1),
            progress_callback=progress_callback,
        ):
            excerpt = re.sub(r"\s+", " ", str(page.get("text", "")).strip())[:1200]
            viewed_pages.append(
                {
                    "url": str(page.get("url", "")).strip(),
                    "content_type": str(page.get("content_type", "")).strip(),
                    "excerpt": excerpt,
                    **({"error": str(page.get("error", "")).strip()} if page.get("error") else {}),
                }
            )

    return {
        "provider": provider,
        "providers_tried": attempts,
        "results": search_results[:bounded_num],
        "viewed_pages": viewed_pages,
        "raw": raw_payload or {},
        "instant_answer": raw_payload.get("instant_answer", {}) if isinstance(raw_payload, dict) else {},
        "query_plan": raw_payload.get("query_plan", []) if isinstance(raw_payload, dict) else [],
        "error": "" if search_results else "; ".join(errors),
    }


def openalex_search(query: str, per_page: int = 10, filters: str | None = None) -> dict:
    params = {
        "search": query,
        "per-page": per_page,
    }
    if filters:
        params["filter"] = filters
    request = Request(
        f"{OPENALEX_API_URL}?{urlencode(params)}",
        headers={
            "User-Agent": os.getenv(
                "RESEARCH_USER_AGENT",
                "multi-agent-research-bot/1.0 (+https://localhost)",
            )
        },
    )
    with urlopen(request, timeout=30) as response:
        return json.loads(response.read().decode("utf-8"))


def arxiv_search(query: str, max_results: int = 10, sort_by: str = "relevance") -> list[dict]:
    """Search arXiv for academic papers using the public Atom feed API (no API key required)."""
    params = {
        "search_query": f"all:{query}",
        "start": 0,
        "max_results": max(1, min(50, max_results)),
        "sortBy": sort_by,
        "sortOrder": "descending",
    }
    request = Request(
        f"{ARXIV_API_URL}?{urlencode(params)}",
        headers={"User-Agent": os.getenv("RESEARCH_USER_AGENT", "multi-agent-research-bot/1.0 (+https://localhost)")},
    )
    with urlopen(request, timeout=30) as response:
        xml_data = response.read().decode("utf-8")
    root = ET.fromstring(xml_data)
    ns = {
        "atom": "http://www.w3.org/2005/Atom",
        "arxiv": "http://arxiv.org/schemas/atom",
    }
    entries = []
    for entry in root.findall("atom:entry", ns):
        title = (entry.findtext("atom:title", "", ns) or "").strip().replace("\n", " ")
        summary = (entry.findtext("atom:summary", "", ns) or "").strip().replace("\n", " ")
        arxiv_url = (entry.findtext("atom:id", "", ns) or "").strip()
        authors = [
            (a.findtext("atom:name", "", ns) or "").strip()
            for a in entry.findall("atom:author", ns)
        ]
        published = (entry.findtext("atom:published", "", ns) or "").strip()
        categories = [c.attrib.get("term", "") for c in entry.findall("atom:category", ns)]
        # Extract year from published date (e.g. "2023-01-15T00:00:00Z" → "2023")
        year = published[:4] if published and len(published) >= 4 else "Unknown"
        # Extract PDF link from alternate links in the entry
        pdf_url = ""
        for link in entry.findall("atom:link", ns):
            if link.attrib.get("type") == "application/pdf" or link.attrib.get("title") == "pdf":
                pdf_url = link.attrib.get("href", "")
                break
        if not pdf_url and arxiv_url and "/abs/" in arxiv_url:
            # Derive PDF URL safely: http://arxiv.org/abs/2301.00001v1 → http://arxiv.org/pdf/2301.00001v1
            # Simply swap /abs/ for /pdf/ — do NOT strip version suffix (v1, v2, etc.)
            pdf_url = arxiv_url.replace("/abs/", "/pdf/", 1)
        entries.append(
            {
                "title": title,
                # Canonical field names required by downstream consumers
                "abstract": summary[:800],
                "year": year,
                "pdf_url": pdf_url,
                # Legacy aliases kept for backward compatibility
                "summary": summary[:800],
                "published": published,
                "url": arxiv_url,
                "authors": [a for a in authors if a][:6],
                "categories": [c for c in categories if c][:5],
                "source": "arxiv",
            }
        )
    return entries


def reddit_search(
    query: str,
    subreddit: str = "",
    sort: str = "relevance",
    limit: int = 10,
) -> list[dict]:
    """Search Reddit posts using the public JSON search API (no authentication required)."""
    sub_raw = str(subreddit or "").strip()
    # Use removeprefix to avoid lstrip's char-set semantics (lstrip("r/") corrupts names like "reactjs")
    subreddit = sub_raw.removeprefix("r/").removeprefix("/r/").strip()
    base = f"{REDDIT_BASE_URL}/r/{subreddit}" if subreddit else REDDIT_BASE_URL
    params = {
        "q": query,
        "sort": sort,
        "limit": max(1, min(100, limit)),
        "type": "link",
        "raw_json": 1,
    }
    url = f"{base}/search.json?{urlencode(params)}"
    request = Request(
        url,
        headers={"User-Agent": os.getenv("RESEARCH_USER_AGENT", "multi-agent-research-bot/1.0 (+https://localhost)")},
    )
    with urlopen(request, timeout=30) as response:
        data = json.loads(response.read().decode("utf-8"))
    posts = []
    for child in data.get("data", {}).get("children", []) or []:
        item = child.get("data", {})
        if not isinstance(item, dict):
            continue
        title = str(item.get("title") or "").strip()
        selftext = str(item.get("selftext") or "").strip()[:800]
        permalink = str(item.get("permalink") or "").strip()
        post_url = f"https://www.reddit.com{permalink}" if permalink.startswith("/") else permalink
        subreddit_name = str(item.get("subreddit") or "").strip()
        score = int(item.get("score") or 0)
        num_comments = int(item.get("num_comments") or 0)
        posts.append(
            {
                "title": title,
                "text": selftext,
                "url": post_url,
                "subreddit": subreddit_name,
                "score": score,
                "num_comments": num_comments,
                "source": "reddit",
            }
        )
    return posts


def search_result_urls(payload: dict) -> list[str]:
    urls = []
    for result in payload.get("organic_results", []):
        link = result.get("link")
        if link and link not in urls:
            urls.append(link)
    return urls


def summarize_pages(pages: list[dict], objective: str, audience: str) -> str:
    prompt = f"""
You are a web research summarizer.

Objective:
{objective}

Audience:
{audience}

Pages:
{json.dumps(pages, indent=2, ensure_ascii=False)}

Produce a concise but evidence-oriented summary. Mention important facts, conflicts, and missing pieces.
"""
    return llm_text(prompt)


def _clean_whitespace(value: str) -> str:
    return re.sub(r"[ \t]+", " ", (value or "")).strip()


def _text_structure_metadata(text: str) -> dict:
    lines = [line.strip() for line in str(text or "").splitlines() if line.strip()]
    heading_count = sum(
        1
        for line in lines
        if line.startswith("#")
        or (len(line) <= 120 and line == line.upper() and re.search(r"[A-Z]", line))
        or line.endswith(":")
    )
    return {
        "line_count": len(lines),
        "heading_count": heading_count,
        "character_count": len(str(text or "")),
    }


def _classify_extract_error(exc: Exception) -> tuple[str, bool]:
    message = str(exc or "").lower()
    if "encrypted" in message or "password" in message:
        return "encrypted", False
    if "ocr mode" in message or "ocr" in message:
        return "ocr_required", True
    if "not found" in message:
        return "missing", False
    if "unsupported" in message:
        return "unsupported", False
    if "badzipfile" in message or "file is not a zip file" in message or "cannot parse" in message:
        return "corrupt", False
    if "pdf" in message and "stream" in message:
        return "corrupt", False
    return "extract_failed", False


def _extract_printable_chunks(data: bytes, *, min_length: int = 6, max_chunks: int = 3000) -> str:
    decoded = data.decode("latin-1", errors="ignore").replace("\x00", " ")
    pattern = rf"[^\x00-\x1F]{{{max(1, min_length)},}}"
    chunks = [_clean_whitespace(item) for item in re.findall(pattern, decoded)]
    filtered = [item for item in chunks if len(item) >= min_length and not item.startswith("PK")]
    return "\n".join(filtered[:max_chunks])


def _run_command_capture_text(command: list[str], timeout: int = 60) -> str:
    executable = command[0]
    if not shutil.which(executable):
        return ""
    try:
        completed = subprocess.run(
            command,
            capture_output=True,
            text=True,
            check=False,
            timeout=timeout,
        )
    except Exception:
        return ""
    if completed.returncode != 0:
        return ""
    return (completed.stdout or "").strip()


def _extract_doc_text(path: Path) -> tuple[str, dict]:
    for command in (["antiword", str(path)], ["catdoc", str(path)]):
        text = _run_command_capture_text(command)
        if text:
            return text, {"type": "doc", "reader": command[0]}
    fallback = _extract_printable_chunks(path.read_bytes())
    if fallback:
        return fallback, {"type": "doc", "reader": "binary_strings"}
    raise ValueError(f"Unable to read legacy DOC file: {path}")


def _extract_docx_xml_text(path: Path) -> tuple[str, dict]:
    with zipfile.ZipFile(path) as archive:
        word_parts = [
            "word/document.xml",
            *sorted(
                item
                for item in archive.namelist()
                if item.startswith("word/header") and item.endswith(".xml")
            ),
            *sorted(
                item
                for item in archive.namelist()
                if item.startswith("word/footer") and item.endswith(".xml")
            ),
            "word/footnotes.xml",
            "word/endnotes.xml",
        ]
        seen: set[str] = set()
        lines: list[str] = []
        table_count = 0
        for part in word_parts:
            if part in seen or part not in archive.namelist():
                continue
            seen.add(part)
            root = ET.fromstring(archive.read(part))
            table_count += len(root.findall(".//{*}tbl"))
            paragraphs = root.findall(".//{*}p") or [root]
            for paragraph in paragraphs:
                pieces = [
                    node.text.strip()
                    for node in paragraph.findall(".//{*}t")
                    if node.text and node.text.strip()
                ]
                if pieces:
                    lines.append(" ".join(pieces))
        text = "\n".join(lines).strip()
        return text, {"type": "docx", "reader": "zipxml", "paragraphs": len(lines), "table_count": table_count}


def _extract_docx_text(path: Path) -> tuple[str, dict]:
    errors: list[str] = []
    try:
        from docx import Document

        doc = Document(str(path))
        text = "\n".join(paragraph.text for paragraph in doc.paragraphs if paragraph.text)
        if text.strip():
            return text, {
                "type": "docx",
                "reader": "python-docx",
                "paragraphs": len(doc.paragraphs),
                "table_count": len(getattr(doc, "tables", []) or []),
            }
        errors.append("python-docx returned no text")
    except Exception as exc:
        errors.append(str(exc))

    try:
        text, metadata = _extract_docx_xml_text(path)
        if text.strip():
            if errors:
                metadata = {**metadata, "primary_extract_error": errors[0]}
            return text, metadata
        errors.append("zipxml returned no text")
    except Exception as exc:
        errors.append(str(exc))

    fallback = _extract_printable_chunks(path.read_bytes())
    if fallback:
        metadata = {"type": "docx", "reader": "binary_strings"}
        if errors:
            metadata["primary_extract_error"] = errors[0]
        return fallback, metadata
    raise ValueError(f"Unable to read DOCX file: {path}. attempts={' | '.join(errors) if errors else 'none'}")


def _extract_xlsx_xml_text(path: Path) -> tuple[str, dict]:
    with zipfile.ZipFile(path) as archive:
        shared_strings = []
        if "xl/sharedStrings.xml" in archive.namelist():
            root = ET.fromstring(archive.read("xl/sharedStrings.xml"))
            shared_strings = [node.text.strip() for node in root.findall(".//{*}t") if node.text and node.text.strip()]

        sheet_paths = sorted(
            item
            for item in archive.namelist()
            if item.startswith("xl/worksheets/") and item.endswith(".xml")
        )

        lines = []
        sheet_names: list[str] = []
        non_empty_rows = 0
        for sheet_index, sheet_path in enumerate(sheet_paths, start=1):
            sheet_name = Path(sheet_path).name
            sheet_names.append(sheet_name)
            lines.append(f"[Sheet {sheet_index}] {sheet_name}")
            root = ET.fromstring(archive.read(sheet_path))
            for row in root.findall(".//{*}row")[:600]:
                row_values = []
                for cell in row.findall("{*}c"):
                    cell_type = cell.attrib.get("t", "")
                    text_value = ""
                    if cell_type == "inlineStr":
                        text_value = " ".join(
                            item.text.strip()
                            for item in cell.findall(".//{*}t")
                            if item.text and item.text.strip()
                        )
                    else:
                        value_node = cell.find("{*}v")
                        if value_node is not None and value_node.text:
                            if cell_type == "s":
                                try:
                                    idx = int(value_node.text)
                                    text_value = shared_strings[idx] if 0 <= idx < len(shared_strings) else value_node.text
                                except Exception:
                                    text_value = value_node.text
                            else:
                                text_value = value_node.text
                    cleaned = _clean_whitespace(text_value)
                    if cleaned:
                        row_values.append(cleaned)
                if row_values:
                    non_empty_rows += 1
                    lines.append(", ".join(row_values))
        text = "\n".join(lines).strip()
        return text, {"type": "xlsx", "reader": "zipxml", "sheets": len(sheet_paths), "sheet_names": sheet_names, "non_empty_rows": non_empty_rows}


def _extract_xlsx_text(path: Path) -> tuple[str, dict]:
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(str(path), data_only=True, read_only=True)
        lines = []
        sheet_names: list[str] = []
        non_empty_rows = 0
        for sheet in workbook.worksheets:
            sheet_names.append(sheet.title)
            lines.append(f"[Sheet] {sheet.title}")
            for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                if row_index > 600:
                    break
                values = [_clean_whitespace(str(value)) for value in row if value is not None and _clean_whitespace(str(value))]
                if values:
                    non_empty_rows += 1
                    lines.append(", ".join(values))
        text = "\n".join(lines).strip()
        return text, {"type": "xlsx", "reader": "openpyxl", "sheets": len(workbook.worksheets), "sheet_names": sheet_names, "non_empty_rows": non_empty_rows}
    except Exception as exc:
        try:
            text, metadata = _extract_xlsx_xml_text(path)
            if text.strip():
                return text, metadata
        except Exception:
            pass
        raise exc


def _extract_xls_text(path: Path) -> tuple[str, dict]:
    try:
        import xlrd

        workbook = xlrd.open_workbook(str(path))
        lines = []
        for sheet in workbook.sheets():
            lines.append(f"[Sheet] {sheet.name}")
            for row_index in range(min(sheet.nrows, 600)):
                values = []
                for col_index in range(sheet.ncols):
                    value = _clean_whitespace(str(sheet.cell_value(row_index, col_index)))
                    if value:
                        values.append(value)
                if values:
                    lines.append(", ".join(values))
        text = "\n".join(lines).strip()
        return text, {"type": "xls", "reader": "xlrd", "sheets": workbook.nsheets}
    except Exception:
        text = _run_command_capture_text(["xls2csv", str(path)])
        if text:
            return text, {"type": "xls", "reader": "xls2csv"}
        fallback = _extract_printable_chunks(path.read_bytes())
        if fallback:
            return fallback, {"type": "xls", "reader": "binary_strings"}
        raise ValueError(f"Unable to read legacy XLS file: {path}")


def _extract_pptx_xml_text(path: Path) -> tuple[str, dict]:
    with zipfile.ZipFile(path) as archive:
        slide_paths = sorted(
            item
            for item in archive.namelist()
            if item.startswith("ppt/slides/slide") and item.endswith(".xml")
        )
        lines = []
        notes_count = 0
        for index, slide_path in enumerate(slide_paths, start=1):
            root = ET.fromstring(archive.read(slide_path))
            parts = [node.text.strip() for node in root.findall(".//{*}t") if node.text and node.text.strip()]
            lines.append(f"[Slide {index}]")
            if parts:
                lines.append(" ".join(parts))
            notes_path = f"ppt/notesSlides/notesSlide{index}.xml"
            if notes_path in archive.namelist():
                notes_root = ET.fromstring(archive.read(notes_path))
                notes_parts = [node.text.strip() for node in notes_root.findall(".//{*}t") if node.text and node.text.strip()]
                if notes_parts:
                    notes_count += 1
                    lines.append(f"[Slide {index} Notes]")
                    lines.append(" ".join(notes_parts))
        text = "\n".join(lines).strip()
        return text, {"type": "pptx", "reader": "zipxml", "slides": len(slide_paths), "notes_count": notes_count}


def _extract_pptx_text(path: Path) -> tuple[str, dict]:
    try:
        from pptx import Presentation

        presentation = Presentation(str(path))
        lines = []
        notes_count = 0
        table_count = 0
        for slide_index, slide in enumerate(presentation.slides, start=1):
            parts = []
            for shape in slide.shapes:
                if getattr(shape, "has_table", False):
                    table_count += 1
                if getattr(shape, "has_text_frame", False) and shape.text:
                    cleaned = _clean_whitespace(shape.text)
                    if cleaned:
                        parts.append(cleaned)
            lines.append(f"[Slide {slide_index}]")
            if parts:
                lines.append(" ".join(parts))
            try:
                notes_frame = slide.notes_slide.notes_text_frame
                notes_text = _clean_whitespace(notes_frame.text if notes_frame is not None else "")
                if notes_text:
                    notes_count += 1
                    lines.append(f"[Slide {slide_index} Notes]")
                    lines.append(notes_text)
            except Exception:
                pass
        text = "\n".join(lines).strip()
        return text, {"type": "pptx", "reader": "python-pptx", "slides": len(presentation.slides), "notes_count": notes_count, "table_count": table_count}
    except Exception as exc:
        errors = [str(exc)]
        try:
            text, metadata = _extract_pptx_xml_text(path)
            if text.strip():
                metadata = {**metadata, "primary_extract_error": errors[0]}
                return text, metadata
            errors.append("zipxml returned no text")
        except Exception as xml_exc:
            errors.append(str(xml_exc))
        fallback = _extract_printable_chunks(path.read_bytes())
        if fallback:
            return fallback, {"type": "pptx", "reader": "binary_strings", "primary_extract_error": errors[0]}
        raise ValueError(f"Unable to read PPTX file: {path}. attempts={' | '.join(errors)}")


def _extract_ppt_text(path: Path) -> tuple[str, dict]:
    text = _run_command_capture_text(["catppt", str(path)])
    if text:
        return text, {"type": "ppt", "reader": "catppt"}
    fallback = _extract_printable_chunks(path.read_bytes())
    if fallback:
        return fallback, {"type": "ppt", "reader": "binary_strings"}
    raise ValueError(f"Unable to read legacy PPT file: {path}")


def _extract_pdf_text(path: Path) -> tuple[str, dict]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    if getattr(reader, "is_encrypted", False):
        try:
            decrypt_status = reader.decrypt("")
        except Exception as exc:
            raise ValueError(f"Encrypted PDF requires password: {path}") from exc
        if not decrypt_status:
            raise ValueError(f"Encrypted PDF requires password: {path}")
    text = "\n".join((page.extract_text() or "") for page in reader.pages)
    return text, {"type": "pdf", "reader": "pypdf", "pages": len(reader.pages)}


def _extract_excel_text_via_pdf_fallback(
    path: Path,
    *,
    source_type: str,
    original_error: Exception | None = None,
) -> tuple[str, dict]:
    converter_binaries = ("soffice", "libreoffice")
    attempts: list[str] = []
    errors: list[str] = []

    for binary in converter_binaries:
        if not shutil.which(binary):
            attempts.append(f"{binary}:not_installed")
            continue

        with tempfile.TemporaryDirectory(prefix="kendr_excel_pdf_") as temp_dir:
            command = [binary, "--headless", "--convert-to", "pdf", "--outdir", temp_dir, str(path)]
            try:
                completed = subprocess.run(
                    command,
                    capture_output=True,
                    text=True,
                    check=False,
                    timeout=180,
                )
            except Exception as exc:
                attempts.append(f"{binary}:exception")
                errors.append(f"{binary} execution failed: {exc}")
                continue

            if completed.returncode != 0:
                attempts.append(f"{binary}:non_zero_exit")
                stderr_text = (completed.stderr or "").strip()
                stdout_text = (completed.stdout or "").strip()
                errors.append(
                    f"{binary} returned code {completed.returncode}. stderr={stderr_text[:240]} stdout={stdout_text[:240]}"
                )
                continue

            attempts.append(f"{binary}:ok")
            expected_pdf = Path(temp_dir) / f"{path.stem}.pdf"
            pdf_candidates = sorted(Path(temp_dir).glob("*.pdf"))
            if expected_pdf.exists():
                pdf_candidates = [expected_pdf] + [candidate for candidate in pdf_candidates if candidate != expected_pdf]

            if not pdf_candidates:
                errors.append(f"{binary} succeeded but did not produce a PDF for {path.name}.")
                continue

            for pdf_path in pdf_candidates:
                try:
                    text, pdf_meta = _extract_pdf_text(pdf_path)
                except Exception as exc:
                    errors.append(f"Failed to parse fallback PDF {pdf_path.name}: {exc}")
                    continue
                if not text.strip():
                    errors.append(f"Fallback PDF {pdf_path.name} had no extractable text.")
                    continue
                metadata = {
                    "type": source_type,
                    "reader": f"{binary}_pdf_fallback",
                    "fallback_source": "pdf_conversion",
                    "fallback_pdf_file": pdf_path.name,
                    "fallback_pdf_pages": int(pdf_meta.get("pages", 0) or 0),
                }
                if original_error is not None:
                    metadata["primary_extract_error"] = str(original_error)
                return text, metadata

    reason = "; ".join(errors) if errors else "no available converter or conversion output"
    attempted = ", ".join(attempts) if attempts else "none"
    raise ValueError(
        f"Excel PDF fallback failed for {path.name}. attempts={attempted}. reason={reason}"
    )


def parse_document(path_str: str, *, ocr_images: bool = False, ocr_instruction: str | None = None) -> dict:
    path = Path(path_str)
    if not path.exists() or not path.is_file():
        raise ValueError(f"File not found for ingestion: {path}")

    suffix = path.suffix.lower()
    if suffix in TEXT_LIKE_SOURCE_EXTENSIONS:
        text = path.read_text(encoding="utf-8", errors="ignore")
        if suffix in {".html", ".htm"}:
            text = html_to_text(text)
        metadata = {"type": suffix.lstrip("."), "reader": "text"}
        metadata.update(_text_structure_metadata(text))
        return {"path": str(path), "text": text, "metadata": metadata}
    if suffix == ".csv":
        with open(path, "r", encoding="utf-8", errors="ignore", newline="") as handle:
            reader = csv.reader(handle)
            rows = list(reader)
        text = "\n".join(", ".join(row) for row in rows[:200])
        metadata = {"type": "csv", "reader": "csv", "rows": len(rows), "columns": max((len(row) for row in rows), default=0)}
        metadata.update(_text_structure_metadata(text))
        return {"path": str(path), "text": text, "metadata": metadata}
    if suffix == ".pdf":
        text, metadata = _extract_pdf_text(path)
        return {"path": str(path), "text": text, "metadata": metadata}
    if suffix == ".docx":
        text, metadata = _extract_docx_text(path)
        return {"path": str(path), "text": text, "metadata": metadata}
    if suffix == ".doc":
        text, metadata = _extract_doc_text(path)
        return {"path": str(path), "text": text, "metadata": metadata}
    if suffix in {".xlsx", ".xlsm"}:
        try:
            text, metadata = _extract_xlsx_text(path)
        except Exception as exc:
            text, metadata = _extract_excel_text_via_pdf_fallback(path, source_type=suffix.lstrip("."), original_error=exc)
        return {"path": str(path), "text": text, "metadata": metadata}
    if suffix == ".xls":
        try:
            text, metadata = _extract_xls_text(path)
        except Exception as exc:
            text, metadata = _extract_excel_text_via_pdf_fallback(path, source_type="xls", original_error=exc)
        return {"path": str(path), "text": text, "metadata": metadata}
    if suffix in {".pptx", ".pptm"}:
        text, metadata = _extract_pptx_text(path)
        return {"path": str(path), "text": text, "metadata": metadata}
    if suffix == ".ppt":
        text, metadata = _extract_ppt_text(path)
        return {"path": str(path), "text": text, "metadata": metadata}
    if suffix in IMAGE_FILE_EXTENSIONS:
        if not ocr_images:
            raise ValueError(f"Image ingestion requires OCR mode for file: {path}")
        result = openai_ocr_image(str(path), instruction=ocr_instruction)
        return {
            "path": str(path),
            "text": result.get("text", ""),
            "metadata": {
                "type": "image",
                "reader": "openai_ocr",
                "ocr_characters": len(result.get("text", "") or ""),
                **_text_structure_metadata(result.get("text", "") or ""),
            },
        }
    raise ValueError(f"Unsupported document type for ingestion: {path}")


def _document_parse_bucket(path_str: str, *, ocr_images: bool) -> str:
    suffix = Path(path_str).suffix.lower()
    if suffix in IMAGE_FILE_EXTENSIONS and ocr_images:
        return "ocr"
    if suffix in {".pdf", ".doc", ".docx", ".xls", ".xlsx", ".xlsm", ".ppt", ".pptx", ".pptm"}:
        return "heavy"
    return "light"


def _document_parse_error(path_str: str, exc: Exception) -> dict:
    resolved_path = Path(path_str)
    error_kind, recoverable = _classify_extract_error(exc)
    return {
        "path": str(resolved_path),
        "text": "",
        "metadata": {
            "type": resolved_path.suffix.lstrip(".") or "unknown",
            "error": str(exc),
            "error_kind": error_kind,
            "recoverable": recoverable,
        },
    }


def parse_documents(
    paths: list[str],
    *,
    continue_on_error: bool = False,
    ocr_images: bool = False,
    ocr_instruction: str | None = None,
    progress_callback=None,
) -> list[dict]:
    ordered_paths = [str(path or "").strip() for path in paths if str(path or "").strip()]
    if not ordered_paths:
        return []

    if len(ordered_paths) == 1:
        path = ordered_paths[0]
        if progress_callback:
            try:
                progress_callback(path, "started", None, 1, 1)
            except Exception:
                pass
        try:
            payload = parse_document(path, ocr_images=ocr_images, ocr_instruction=ocr_instruction)
            if progress_callback:
                try:
                    progress_callback(path, "completed", payload, 1, 1)
                except Exception:
                    pass
            return [payload]
        except Exception as exc:
            if not continue_on_error:
                raise
            payload = _document_parse_error(path, exc)
            if progress_callback:
                try:
                    progress_callback(path, "failed", payload, 1, 1)
                except Exception:
                    pass
            return [payload]

    bucket_limits = {
        "light": _document_worker_count(len(ordered_paths)),
        "heavy": _heavy_document_worker_count(),
        "ocr": _ocr_document_worker_count(),
    }
    semaphore_by_bucket = {
        bucket: threading.BoundedSemaphore(limit)
        for bucket, limit in bucket_limits.items()
    }
    results: list[dict | None] = [None] * len(ordered_paths)
    max_workers = min(len(ordered_paths), sum(bucket_limits.values()))

    def _parse_one(index: int, path: str) -> tuple[int, dict]:
        bucket = _document_parse_bucket(path, ocr_images=ocr_images)
        with semaphore_by_bucket[bucket]:
            if progress_callback:
                try:
                    progress_callback(path, "started", None, index + 1, len(ordered_paths))
                except Exception:
                    pass
            try:
                payload = parse_document(path, ocr_images=ocr_images, ocr_instruction=ocr_instruction)
                if progress_callback:
                    try:
                        progress_callback(path, "completed", payload, index + 1, len(ordered_paths))
                    except Exception:
                        pass
                return index, payload
            except Exception as exc:
                if not continue_on_error:
                    raise
                payload = _document_parse_error(path, exc)
                if progress_callback:
                    try:
                        progress_callback(path, "failed", payload, index + 1, len(ordered_paths))
                    except Exception:
                        pass
                return index, payload

    with ThreadPoolExecutor(max_workers=max_workers) as pool:
        futures = {
            pool.submit(_parse_one, index, path): index
            for index, path in enumerate(ordered_paths)
        }
        try:
            for future in as_completed(futures):
                index, payload = future.result()
                results[index] = payload
        except Exception:
            for future in futures:
                future.cancel()
            raise

    return [item for item in results if isinstance(item, dict)]


def openai_ocr_image(path_str: str, instruction: str | None = None) -> dict:
    return _run_openai_compat_image_prompt(
        path_str,
        instruction or "Extract all visible text and key structured details from this image. Return plain text.",
        agent_name="ocr_agent",
    )


def openai_analyze_image(path_str: str, instruction: str | None = None) -> dict:
    result = _run_openai_compat_image_prompt(
        path_str,
        instruction
        or (
            "Analyze this image and extract meaningful information. "
            "Describe the scene, objects, text, layout, visual signals, "
            "data patterns, probable context, and actionable insights if present."
        ),
        agent_name="image_agent",
    )
    return {
        "path": result.get("path", path_str),
        "analysis": result.get("text", ""),
        "raw": result.get("raw", {}),
        "provider": result.get("provider", ""),
        "model": result.get("model", ""),
        "selection_source": result.get("selection_source", ""),
    }


def get_qdrant_client():
    from qdrant_client import QdrantClient

    return QdrantClient(url=DEFAULT_QDRANT_URL)


def ensure_vector_collection(collection_name: str = DEFAULT_QDRANT_COLLECTION, vector_size: int = 1536):
    from tasks.vector_backends import get_vector_backend

    backend = get_vector_backend()
    backend.ensure_collection(collection_name, vector_size=vector_size)
    return backend


def embed_texts(texts: list[str], model: str = DEFAULT_EMBEDDING_MODEL) -> list[list[float]]:
    if not texts:
        return []
    client = get_openai_client()
    response = client.embeddings.create(model=model, input=texts)
    return [item.embedding for item in response.data]


def upsert_memory_records(records: list[dict], collection_name: str = DEFAULT_QDRANT_COLLECTION):
    from tasks.vector_backends import get_vector_backend

    if not records:
        return {"indexed": 0, "collection": collection_name}
    vectors = embed_texts([record["text"] for record in records])
    backend = get_vector_backend()
    backend.ensure_collection(collection_name, vector_size=len(vectors[0]) if vectors else 1536)
    return backend.upsert(collection_name, records, vectors)


def search_memory(query: str, top_k: int = 5, collection_name: str = DEFAULT_QDRANT_COLLECTION) -> list[dict]:
    from tasks.vector_backends import get_vector_backend

    backend = get_vector_backend()
    backend.ensure_collection(collection_name)
    query_vector = embed_texts([query])[0]
    return backend.search(collection_name, query_vector, top_k=top_k)


def build_evidence_bundle(state: dict) -> dict:
    return {
        "user_query": state.get("user_query", ""),
        "current_objective": state.get("current_objective", ""),
        "search_summary": state.get("search_summary", ""),
        "research_result": state.get("research_result", ""),
        "web_crawl_summary": state.get("web_crawl_summary", ""),
        "document_summary": state.get("document_summary", ""),
        "ocr_summary": state.get("ocr_summary", ""),
        "entity_resolution": state.get("entity_resolution", {}),
        "knowledge_graph": state.get("knowledge_graph", {}),
        "timeline": state.get("timeline", []),
        "people_profile": state.get("people_profile", {}),
        "company_profile": state.get("company_profile", {}),
        "relationship_map": state.get("relationship_map", {}),
        "compliance_risk_report": state.get("compliance_risk_report", {}),
        "verification_report": state.get("verification_report", {}),
        "citations": state.get("citations", []),
        "structured_facts": state.get("structured_facts", []),
        "agent_history": state.get("agent_history", [])[-10:],
    }


def evidence_text(evidence: dict) -> str:
    return json.dumps(evidence, indent=2, ensure_ascii=False)
